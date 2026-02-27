"""Tests for tools/files.py - file operations."""

import json
import logging
from unittest.mock import MagicMock, patch

import pytest
from docx import Document as DocxDocument
from openpyxl import Workbook
from pptx import Presentation

from services.vault import clear_pending_previews
from tools.readers import handle_audio, handle_image, handle_office
from tools.files import (
    _embed_cache,
    _expand_embeds,
    _extract_block,
    _extract_headings,
    _merge_bodies,
    _merge_frontmatter,
    _split_blocks,
    _split_frontmatter_body,
    batch_merge_files,
    batch_move_files,
    create_file,
    get_note_info,
    merge_files,
    move_file,
    read_file,
)


class TestReadFile:
    """Tests for read_file tool."""

    def test_read_existing_file(self, vault_config):
        """Should read content of existing file."""
        result = json.loads(read_file("note1.md"))
        assert result["success"] is True
        assert "# Note 1" in result["content"]
        assert "wikilink" in result["content"]

    def test_read_file_not_found(self, vault_config):
        """Should return error for missing file."""
        result = json.loads(read_file("nonexistent.md"))
        assert result["success"] is False
        assert "not found" in result["error"].lower()

    def test_read_file_in_subdirectory(self, vault_config):
        """Should read file in subdirectory."""
        result = json.loads(read_file("projects/project1.md"))
        assert result["success"] is True
        assert "# Project 1" in result["content"]

    def test_short_file_no_markers(self, vault_config):
        """Short files should be returned in full with no markers."""
        result = json.loads(read_file("note3.md"))
        assert result["success"] is True
        assert "[... truncated" not in result["content"]
        assert "[Continuing from" not in result["content"]
        assert "# Note 3" in result["content"]

    def test_long_file_truncated_with_marker(self, vault_config):
        """Files longer than length should have a continuation marker."""
        long_content = "# Long Note\n\n" + "x" * 50000
        (vault_config / "long.md").write_text(long_content)
        result = json.loads(read_file("long.md"))
        assert result["success"] is True
        assert result["content"].startswith("# Long Note")
        assert "[... truncated at char 30000 of" in result["content"]
        assert "Use offset=30000 to read more." in result["content"]

    def test_offset_pagination(self, vault_config):
        """Reading with offset should show continuation header and may show truncation marker."""
        long_content = "A" * 100000
        (vault_config / "long.md").write_text(long_content)
        result = json.loads(read_file("long.md", offset=30000))
        assert result["success"] is True
        assert "[Continuing from char 30000 of 100000]" in result["content"]
        assert "[... truncated at char 60000 of 100000" in result["content"]

    def test_offset_final_chunk(self, vault_config):
        """Reading the last chunk should have no truncation marker."""
        long_content = "B" * 50000
        (vault_config / "long.md").write_text(long_content)
        result = json.loads(read_file("long.md", offset=30000))
        assert result["success"] is True
        assert "[Continuing from char 30000 of 50000]" in result["content"]
        assert "[... truncated" not in result["content"]

    def test_offset_past_end(self, vault_config):
        """Offset past end of file should return error."""
        result = json.loads(read_file("note3.md", offset=99999))
        assert result["success"] is False
        assert "offset" in result["error"].lower()

    def test_read_file_utf8_encoding(self, vault_config):
        """Should handle UTF-8 content including non-ASCII characters."""
        content = "# Café\n\nRésumé with émojis: 🎉"
        (vault_config / "utf8.md").write_text(content, encoding="utf-8")
        result = json.loads(read_file("utf8.md"))
        assert result["success"] is True
        assert "Café" in result["content"]
        assert "🎉" in result["content"]

    def test_custom_length(self, vault_config):
        """Custom length parameter should control chunk size."""
        content = "C" * 500
        (vault_config / "custom.md").write_text(content)
        result = json.loads(read_file("custom.md", length=100))
        assert result["success"] is True
        assert "[... truncated at char 100 of 500" in result["content"]
        # Content before the truncation marker should be 100 chars
        before_marker = result["content"].split("\n\n[... truncated")[0]
        assert len(before_marker) == 100

    def test_read_file_expands_embeds(self, vault_config):
        """read_file on a .md file with embeds should auto-expand them."""
        (vault_config / "parent.md").write_text(
            "# Parent\n\nSee: ![[note3]]\n\nEnd.\n"
        )
        result = json.loads(read_file("parent.md"))
        assert result["success"] is True
        assert "> [Embedded: note3]" in result["content"]
        assert "> # Note 3" in result["content"]
        assert "![[note3]]" not in result["content"]

    def test_read_file_embeds_pagination(self, vault_config):
        """Pagination offsets apply to expanded content."""
        body = "x" * 100
        (vault_config / "embedded_target.md").write_text(f"# Target\n\n{body}\n")
        (vault_config / "paginate.md").write_text(
            "# Start\n\n![[embedded_target]]\n\n" + "y" * 50000
        )
        result = json.loads(read_file("paginate.md"))
        assert result["success"] is True
        assert "> [Embedded: embedded_target]" in result["content"]
        assert "[... truncated" in result["content"]

    def test_read_file_normalizes_nbsp(self, vault_config):
        """Non-breaking spaces in paths are normalized to regular spaces."""
        (vault_config / "my file.md").write_text("# Content\n")
        result = json.loads(read_file("my\xa0file.md"))
        assert result["success"] is True
        assert "# Content" in result["content"]

    def test_read_non_md_file_no_expansion(self, vault_config):
        """Non-.md text files should not have embeds expanded."""
        (vault_config / "data.txt").write_text("literal ![[note3]] text")
        result = json.loads(read_file("data.txt"))
        assert result["success"] is True
        assert "![[note3]]" in result["content"]
        assert "> [Embedded:" not in result["content"]


class TestReadFileAudio:
    """Tests for read_file dispatching to audio handler."""

    def test_audio_no_api_key(self, vault_config, monkeypatch):
        """Audio files require FIREWORKS_API_KEY."""
        monkeypatch.delenv("FIREWORKS_API_KEY", raising=False)
        audio = vault_config / "Attachments" / "test.m4a"
        audio.write_bytes(b"fake audio")
        result = json.loads(read_file("Attachments/test.m4a"))
        assert result["success"] is False
        assert "FIREWORKS_API_KEY" in result["error"]

    @patch("tools.readers.OpenAI")
    def test_audio_successful(self, mock_openai_class, vault_config, monkeypatch):
        """Audio files are transcribed via Whisper."""
        monkeypatch.setenv("FIREWORKS_API_KEY", "test-key")
        audio = vault_config / "Attachments" / "test.m4a"
        audio.write_bytes(b"fake audio")

        mock_client = MagicMock()
        mock_openai_class.return_value = mock_client
        mock_response = MagicMock()
        mock_response.text = "Hello world"
        mock_client.audio.transcriptions.create.return_value = mock_response

        result = json.loads(read_file("Attachments/test.m4a"))
        assert result["success"] is True
        assert result["transcript"] == "Hello world"

    @patch("tools.readers.OpenAI")
    def test_audio_api_error(self, mock_openai_class, vault_config, monkeypatch):
        """API errors are returned gracefully."""
        monkeypatch.setenv("FIREWORKS_API_KEY", "test-key")
        audio = vault_config / "Attachments" / "test.wav"
        audio.write_bytes(b"fake audio")

        mock_client = MagicMock()
        mock_openai_class.return_value = mock_client
        mock_client.audio.transcriptions.create.side_effect = Exception("Rate limit")

        result = json.loads(read_file("Attachments/test.wav"))
        assert result["success"] is False
        assert "Rate limit" in result["error"]

    def test_audio_extensions_dispatched(self, vault_config, monkeypatch):
        """All audio extensions route to the audio handler."""
        monkeypatch.delenv("FIREWORKS_API_KEY", raising=False)
        for ext in [".m4a", ".mp3", ".wav", ".ogg", ".webm"]:
            f = vault_config / "Attachments" / f"test{ext}"
            f.write_bytes(b"audio")
            result = json.loads(read_file(f"Attachments/test{ext}"))
            assert result["success"] is False
            assert "FIREWORKS_API_KEY" in result["error"], f"Extension {ext} not dispatched to audio handler"


class TestReadFileAttachmentsFallback:
    """Tests for binary files resolving via Attachments directory fallback."""

    def test_bare_audio_name_resolves_to_attachments(self, vault_config, monkeypatch):
        """Bare embed filename (no Attachments/ prefix) falls back to Attachments dir."""
        monkeypatch.delenv("FIREWORKS_API_KEY", raising=False)
        audio = vault_config / "Attachments" / "meeting.m4a"
        audio.write_bytes(b"audio data")
        # Call with bare name, no "Attachments/" prefix
        result = json.loads(read_file("meeting.m4a"))
        assert result["success"] is False
        assert "FIREWORKS_API_KEY" in result["error"]  # reached handler, not "not found"

    def test_bare_docx_name_resolves_to_attachments(self, vault_config):
        """Bare .docx filename falls back to Attachments dir."""
        from docx import Document as DocxDocument
        doc = DocxDocument()
        doc.add_paragraph("Test content")
        path = vault_config / "Attachments" / "report.docx"
        doc.save(str(path))
        result = json.loads(read_file("report.docx"))
        assert result["success"] is True
        assert "Test content" in result["content"]

    def test_bare_image_name_resolves_to_attachments(self, vault_config, monkeypatch):
        """Bare image filename falls back to Attachments dir."""
        monkeypatch.delenv("FIREWORKS_API_KEY", raising=False)
        img = vault_config / "Attachments" / "photo.png"
        img.write_bytes(b"image data")
        result = json.loads(read_file("photo.png"))
        assert result["success"] is False
        assert "FIREWORKS_API_KEY" in result["error"]

    def test_explicit_attachments_path_still_works(self, vault_config, monkeypatch):
        """Explicit Attachments/ prefix still resolves directly."""
        monkeypatch.delenv("FIREWORKS_API_KEY", raising=False)
        audio = vault_config / "Attachments" / "test.mp3"
        audio.write_bytes(b"audio")
        result = json.loads(read_file("Attachments/test.mp3"))
        assert result["success"] is False
        assert "FIREWORKS_API_KEY" in result["error"]

    def test_missing_binary_file_returns_error(self, vault_config):
        """Binary file not found in vault root or Attachments returns error."""
        result = json.loads(read_file("nonexistent.docx"))
        assert result["success"] is False
        assert "not found" in result["error"].lower()

    def test_text_file_no_fallback(self, vault_config):
        """Text files do NOT fall back to Attachments — only binary extensions."""
        result = json.loads(read_file("nonexistent.md"))
        assert result["success"] is False
        assert "not found" in result["error"].lower()


class TestReadFileImage:
    """Tests for read_file dispatching to image handler."""

    def test_image_no_api_key(self, vault_config, monkeypatch):
        monkeypatch.delenv("FIREWORKS_API_KEY", raising=False)
        img = vault_config / "Attachments" / "photo.png"
        img.write_bytes(b"\x89PNG fake image")
        result = json.loads(read_file("Attachments/photo.png"))
        assert result["success"] is False
        assert "FIREWORKS_API_KEY" in result["error"]

    @patch("tools.readers.OpenAI")
    def test_image_successful(self, mock_openai_class, vault_config, monkeypatch):
        monkeypatch.setenv("FIREWORKS_API_KEY", "test-key")
        img = vault_config / "Attachments" / "photo.jpg"
        img.write_bytes(b"fake image data")

        mock_client = MagicMock()
        mock_openai_class.return_value = mock_client
        mock_choice = MagicMock()
        mock_choice.message.content = "A photo of a cat"
        mock_response = MagicMock()
        mock_response.choices = [mock_choice]
        mock_client.chat.completions.create.return_value = mock_response

        result = json.loads(read_file("Attachments/photo.jpg"))
        assert result["success"] is True
        assert result["description"] == "A photo of a cat"

    @patch("tools.readers.OpenAI")
    def test_image_api_error(self, mock_openai_class, vault_config, monkeypatch):
        monkeypatch.setenv("FIREWORKS_API_KEY", "test-key")
        img = vault_config / "Attachments" / "photo.webp"
        img.write_bytes(b"fake image")

        mock_client = MagicMock()
        mock_openai_class.return_value = mock_client
        mock_client.chat.completions.create.side_effect = Exception("Model unavailable")

        result = json.loads(read_file("Attachments/photo.webp"))
        assert result["success"] is False
        assert "Model unavailable" in result["error"]

    def test_image_extensions_dispatched(self, vault_config, monkeypatch):
        monkeypatch.delenv("FIREWORKS_API_KEY", raising=False)
        for ext in [".png", ".jpg", ".jpeg", ".gif", ".webp", ".svg"]:
            f = vault_config / "Attachments" / f"test{ext}"
            f.write_bytes(b"img")
            result = json.loads(read_file(f"Attachments/test{ext}"))
            assert result["success"] is False
            assert "FIREWORKS_API_KEY" in result["error"], f"Extension {ext} not dispatched"


class TestCreateFile:
    """Tests for create_file tool."""

    def test_create_simple_file(self, vault_config):
        """Should create a new file."""
        result = json.loads(create_file("new_note.md", "# New Note\n\nContent here."))
        assert result["success"] is True
        assert "new_note.md" in result["path"]

        # Verify file was created
        content = (vault_config / "new_note.md").read_text()
        assert "# New Note" in content

    def test_create_file_with_frontmatter_dict(self, vault_config):
        """Should create file with native dict frontmatter."""
        result = json.loads(create_file(
            "with_fm_dict.md",
            "Content",
            frontmatter={"tags": ["test"], "status": "draft"}
        ))
        assert result["success"] is True

        content = (vault_config / "with_fm_dict.md").read_text()
        assert "---" in content
        assert "tags:" in content
        assert "test" in content

    def test_create_file_with_frontmatter_json_string(self, vault_config):
        """Should create file with JSON string frontmatter for backward compatibility."""
        result = json.loads(create_file(
            "with_fm_json.md",
            "Content",
            frontmatter='{"tags": ["test"], "status": "draft"}'
        ))
        assert result["success"] is True

        content = (vault_config / "with_fm_json.md").read_text()
        assert "---" in content
        assert "tags:" in content
        assert "test" in content

    def test_create_file_already_exists(self, vault_config):
        """Should return error if file exists."""
        result = json.loads(create_file("note1.md", "content"))
        assert result["success"] is False
        assert "already exists" in result["error"].lower()

    def test_create_file_in_new_directory(self, vault_config):
        """Should create parent directories if needed."""
        result = json.loads(create_file("new_dir/nested/note.md", "content"))
        assert result["success"] is True
        assert (vault_config / "new_dir" / "nested" / "note.md").exists()

    def test_create_file_invalid_frontmatter(self, vault_config):
        """Should return error for invalid JSON frontmatter."""
        result = json.loads(create_file("test.md", "content", frontmatter="not valid json"))
        assert result["success"] is False
        assert "invalid frontmatter json" in result["error"].lower()

    def test_create_file_non_object_frontmatter_json(self, vault_config):
        """Should return error when frontmatter JSON is not an object."""
        result = json.loads(create_file("test.md", "content", frontmatter='["not", "an", "object"]'))
        assert result["success"] is False
        assert "expected a json object" in result["error"].lower()


class TestMoveFile:
    """Tests for move_file tool."""

    def test_move_file_success(self, vault_config):
        """Should move file to new location."""
        result = json.loads(move_file("note3.md", "archive/note3.md"))
        assert result["success"] is True
        assert "Moved" in result["message"]

        assert not (vault_config / "note3.md").exists()
        assert (vault_config / "archive" / "note3.md").exists()

    def test_move_file_source_not_found(self, vault_config):
        """Should return error for missing source."""
        result = json.loads(move_file("nonexistent.md", "destination.md"))
        assert result["success"] is False
        assert "not found" in result["error"].lower()

    def test_move_file_destination_exists(self, vault_config):
        """Should return error if destination exists."""
        result = json.loads(move_file("note1.md", "note2.md"))
        assert result["success"] is False
        assert "already exists" in result["error"].lower()

    def test_move_file_same_path(self, vault_config):
        """Should handle moving to same location."""
        result = json.loads(move_file("note1.md", "note1.md"))
        assert result["success"] is True
        assert "Already at destination" in result["message"]



class TestBatchMoveConfirmationGate:
    """Tests for batch_move_files confirmation requirement."""

    def _create_files(self, vault_config, count):
        """Create N test files and return move dicts."""
        (vault_config / "dest").mkdir(exist_ok=True)
        moves = []
        for i in range(count):
            name = f"move_test_{i}.md"
            (vault_config / name).write_text(f"# Note {i}\n")
            moves.append({"source": name, "destination": f"dest/{name}"})
        return moves

    def test_requires_confirmation_over_threshold(self, vault_config):
        """Should return preview when move count exceeds threshold."""
        clear_pending_previews()
        moves = self._create_files(vault_config, 10)
        result = json.loads(batch_move_files(moves=moves))
        assert result["success"] is True
        assert result["confirmation_required"] is True
        assert "10 files" in result["preview_message"]
        assert len(result["files"]) == 10
        # Verify no files were actually moved
        for move in moves:
            assert (vault_config / move["source"]).exists()
            assert not (vault_config / move["destination"]).exists()

    def test_executes_with_confirm_true(self, vault_config):
        """Should execute when confirm=True even over threshold."""
        clear_pending_previews()
        moves = self._create_files(vault_config, 10)
        # First call: store preview
        batch_move_files(moves=moves)
        # Second call: confirm execution
        result = json.loads(batch_move_files(moves=moves, confirm=True))
        assert result["success"] is True
        assert "confirmation_required" not in result
        assert "10 succeeded" in result["message"]

    def test_executes_under_threshold_without_confirm(self, vault_config):
        """Should execute without confirm when move count is at or below threshold."""
        moves = self._create_files(vault_config, 3)
        result = json.loads(batch_move_files(moves=moves))
        assert result["success"] is True
        assert "confirmation_required" not in result
        assert "3 succeeded" in result["message"]

    def test_confirm_true_without_preview_returns_preview(self, vault_config):
        """Passing confirm=True on first call should still return preview."""
        clear_pending_previews()
        moves = self._create_files(vault_config, 10)
        result = json.loads(batch_move_files(moves=moves, confirm=True))
        assert result["success"] is True
        assert result["confirmation_required"] is True
        # No files should be moved
        for move in moves:
            assert (vault_config / move["source"]).exists()
            assert not (vault_config / move["destination"]).exists()

    def test_two_step_confirmation_flow(self, vault_config):
        """Preview then confirm should execute the batch move."""
        clear_pending_previews()
        moves = self._create_files(vault_config, 8)
        # Step 1: preview
        preview = json.loads(batch_move_files(moves=moves))
        assert preview["confirmation_required"] is True
        # Step 2: confirm
        result = json.loads(batch_move_files(moves=moves, confirm=True))
        assert result["success"] is True
        assert "confirmation_required" not in result
        assert "8 succeeded" in result["message"]

    def test_batch_move_preview_has_preview_message(self, vault_config):
        """Preview should include separate preview_message for UI display."""
        clear_pending_previews()
        moves = []
        for i in range(10):
            path = f"move_preview_{i}.md"
            (vault_config / path).write_text(f"---\ntitle: test{i}\n---\n")
            moves.append({"source": path, "destination": f"dest/move_preview_{i}.md"})
        (vault_config / "dest").mkdir(exist_ok=True)
        result = json.loads(batch_move_files(moves=moves))
        assert "preview_message" in result
        assert "This will move 10 files" in result["preview_message"]
        assert "confirm=true" not in result["preview_message"]

    def test_confirmation_is_single_use(self, vault_config):
        """After execution, same confirm=True requires a new preview."""
        clear_pending_previews()
        moves = self._create_files(vault_config, 8)
        # Step 1: preview
        batch_move_files(moves=moves)
        # Step 2: confirm and execute
        batch_move_files(moves=moves, confirm=True)
        # Step 3: confirm again without new preview — should return preview
        # Recreate files since they were moved
        for move in moves:
            (vault_config / move["source"]).write_text("# Recreated\n")
        result = json.loads(batch_move_files(moves=moves, confirm=True))
        assert result["confirmation_required"] is True


class TestSplitFrontmatterBody:
    """Tests for _split_frontmatter_body helper."""

    def test_file_with_frontmatter(self):
        content = "---\ntitle: Test\ntags:\n  - a\n---\n\n# Body\n\nParagraph."
        fm, body = _split_frontmatter_body(content)
        assert fm == {"title": "Test", "tags": ["a"]}
        assert body.strip() == "# Body\n\nParagraph."

    def test_file_without_frontmatter(self):
        content = "# Just a heading\n\nSome text."
        fm, body = _split_frontmatter_body(content)
        assert fm == {}
        assert body == content

    def test_empty_frontmatter(self):
        content = "---\n---\n\nBody text."
        fm, body = _split_frontmatter_body(content)
        assert fm == {}
        assert body.strip() == "Body text."

    def test_frontmatter_with_empty_body(self):
        content = "---\ntitle: Note\n---\n"
        fm, body = _split_frontmatter_body(content)
        assert fm == {"title": "Note"}
        assert body.strip() == ""

    def test_malformed_yaml_treated_as_no_frontmatter(self):
        content = "---\ntitle: [invalid yaml\n---\n\nBody text."
        fm, body = _split_frontmatter_body(content)
        assert fm == {}
        assert body == content

    def test_frontmatter_no_trailing_newline(self):
        """Frontmatter block ending at EOF without trailing newline."""
        content = "---\ntitle: Note\n---"
        fm, body = _split_frontmatter_body(content)
        assert fm == {"title": "Note"}
        assert body == ""

    def test_indented_dashes_in_block_scalar(self):
        """Indented --- inside a YAML block scalar should not close frontmatter."""
        content = "---\ntitle: Test\ndesc: |\n  text\n  ---\n  more\n---\n\nBody."
        fm, body = _split_frontmatter_body(content)
        assert fm["title"] == "Test"
        assert "more" in fm["desc"]
        assert body.strip() == "Body."

    def test_non_dict_yaml_treated_as_no_frontmatter(self):
        """YAML that parses to a list or scalar is not valid frontmatter."""
        content = "---\n- a\n- b\n---\n\nBody text."
        fm, body = _split_frontmatter_body(content)
        assert fm == {}
        assert body == content


class TestMergeFrontmatter:
    """Tests for _merge_frontmatter helper."""

    def test_source_adds_new_fields(self):
        source = {"author": "Alice", "tags": ["draft"]}
        dest = {"title": "Note"}
        merged = _merge_frontmatter(source, dest)
        assert merged == {"title": "Note", "author": "Alice", "tags": ["draft"]}

    def test_destination_wins_scalar_conflict(self):
        source = {"title": "Old Title", "author": "Alice"}
        dest = {"title": "New Title"}
        merged = _merge_frontmatter(source, dest)
        assert merged["title"] == "New Title"
        assert merged["author"] == "Alice"

    def test_list_fields_union_deduped(self):
        source = {"tags": ["a", "b", "c"]}
        dest = {"tags": ["b", "c", "d"]}
        merged = _merge_frontmatter(source, dest)
        assert merged["tags"] == ["b", "c", "d", "a"]

    def test_source_list_dest_scalar_dest_wins(self):
        source = {"status": ["draft", "review"]}
        dest = {"status": "published"}
        merged = _merge_frontmatter(source, dest)
        assert merged["status"] == "published"

    def test_both_empty(self):
        assert _merge_frontmatter({}, {}) == {}

    def test_source_empty(self):
        dest = {"title": "Keep"}
        assert _merge_frontmatter({}, dest) == {"title": "Keep"}

    def test_dest_empty(self):
        source = {"title": "Bring"}
        assert _merge_frontmatter(source, {}) == {"title": "Bring"}

    def test_identical_frontmatter_unchanged(self):
        fm = {"title": "Same", "tags": ["a", "b"]}
        merged = _merge_frontmatter(fm.copy(), fm.copy())
        assert merged == fm

    def test_does_not_mutate_inputs(self):
        """Merging should not modify the original dest_fm dict."""
        source = {"tags": ["a"]}
        dest = {"tags": ["b"]}
        dest_copy = {"tags": ["b"]}
        merged = _merge_frontmatter(source, dest)
        assert dest == dest_copy  # dest unchanged
        assert merged != dest  # merged is different
        assert merged["tags"] == ["b", "a"]

    def test_dict_items_in_list_field(self):
        """List items that are dicts (unhashable) should not crash."""
        source = {"tags": [{"name": "x"}, {"name": "y"}]}
        dest = {"tags": [{"name": "y"}, {"name": "z"}]}
        merged = _merge_frontmatter(source, dest)
        assert {"name": "y"} in merged["tags"]
        assert {"name": "z"} in merged["tags"]
        assert {"name": "x"} in merged["tags"]
        assert len(merged["tags"]) == 3


class TestSplitBlocks:
    """Tests for _split_blocks helper."""

    def test_split_by_headings(self):
        body = "# Intro\n\nParagraph.\n\n## Tasks\n\n- Item 1\n- Item 2\n"
        blocks = _split_blocks(body)
        assert len(blocks) == 2
        assert blocks[0] == ("# Intro", "# Intro\n\nParagraph.\n\n")
        assert blocks[1] == ("## Tasks", "## Tasks\n\n- Item 1\n- Item 2\n")

    def test_content_before_first_heading(self):
        body = "Some intro text.\n\n# Heading\n\nContent.\n"
        blocks = _split_blocks(body)
        assert len(blocks) == 2
        assert blocks[0] == (None, "Some intro text.\n\n")
        assert blocks[1] == ("# Heading", "# Heading\n\nContent.\n")

    def test_no_headings(self):
        body = "Just a paragraph.\n\nAnother paragraph.\n"
        blocks = _split_blocks(body)
        assert len(blocks) == 1
        assert blocks[0] == (None, body)

    def test_empty_body(self):
        blocks = _split_blocks("")
        assert blocks == []

    def test_whitespace_only(self):
        blocks = _split_blocks("  \n\n  \n")
        assert blocks == []


class TestMergeBodies:
    """Tests for _merge_bodies helper."""

    def test_identical_bodies_no_change(self):
        body = "# Tasks\n\n- Item 1\n- Item 2\n"
        merged, stats = _merge_bodies(body, body)
        assert merged == body
        assert stats["blocks_added"] == 0

    def test_source_has_unique_block_under_existing_heading(self):
        source = "# Tasks\n\n- Item 1\n\n# Notes\n\nSource note.\n"
        dest = "# Tasks\n\n- Item 1\n"
        merged, stats = _merge_bodies(source, dest)
        assert "Source note." in merged
        assert "# Notes" in merged
        assert stats["blocks_added"] == 1

    def test_source_unique_block_appended_when_no_heading_match(self):
        source = "# Unrelated\n\nNew stuff.\n"
        dest = "# Tasks\n\n- Item 1\n"
        merged, stats = _merge_bodies(source, dest)
        assert "New stuff." in merged
        assert "# Unrelated" in merged
        assert merged.startswith("# Tasks\n")
        assert stats["blocks_added"] == 1

    def test_duplicate_blocks_not_added(self):
        source = "# Tasks\n\n- Item 1\n\n# Notes\n\nShared note.\n"
        dest = "# Tasks\n\n- Item 1\n\n# Notes\n\nShared note.\n"
        merged, stats = _merge_bodies(source, dest)
        assert merged == dest
        assert stats["blocks_added"] == 0

    def test_partial_overlap(self):
        source = "# Tasks\n\n- Item 1\n\n# Log\n\nEntry A.\n"
        dest = "# Tasks\n\n- Item 1\n\n# Log\n\nEntry B.\n"
        merged, stats = _merge_bodies(source, dest)
        assert "Entry A." in merged
        assert "Entry B." in merged
        assert stats["blocks_added"] == 1

    def test_empty_source_body(self):
        dest = "# Tasks\n\n- Item 1\n"
        merged, stats = _merge_bodies("", dest)
        assert merged == dest
        assert stats["blocks_added"] == 0

    def test_empty_dest_body(self):
        source = "# Tasks\n\n- Item 1\n"
        merged, stats = _merge_bodies(source, "")
        assert "# Tasks" in merged
        assert "- Item 1" in merged
        assert stats["blocks_added"] == 1


class TestMergeFiles:
    """Tests for merge_files tool."""

    def test_identical_files_deletes_source(self, vault_config):
        """Identical files: source deleted, dest unchanged."""
        content = "---\ntitle: Note\n---\n\n# Body\n\nSame content.\n"
        (vault_config / "src.md").write_text(content)
        (vault_config / "dst.md").write_text(content)
        result = json.loads(merge_files("src.md", "dst.md"))
        assert result["success"] is True
        assert not (vault_config / "src.md").exists()
        assert (vault_config / "dst.md").read_text() == content
        assert result["action"] == "identical"

    def test_frontmatter_only_diff(self, vault_config):
        """Source has extra frontmatter fields, bodies identical."""
        src = "---\ntitle: Note\nauthor: Alice\n---\n\n# Body\n\nText.\n"
        dst = "---\ntitle: Note\n---\n\n# Body\n\nText.\n"
        (vault_config / "src.md").write_text(src)
        (vault_config / "dst.md").write_text(dst)
        result = json.loads(merge_files("src.md", "dst.md"))
        assert result["success"] is True
        assert not (vault_config / "src.md").exists()
        merged = (vault_config / "dst.md").read_text()
        assert "author: Alice" in merged
        assert "title: Note" in merged
        assert result["action"] == "frontmatter_merged"

    def test_body_diff_appended(self, vault_config):
        """Source has unique section, merged into dest."""
        src = "# Tasks\n\n- Item 1\n\n# Extra\n\nNew content.\n"
        dst = "# Tasks\n\n- Item 1\n"
        (vault_config / "src.md").write_text(src)
        (vault_config / "dst.md").write_text(dst)
        result = json.loads(merge_files("src.md", "dst.md"))
        assert result["success"] is True
        merged = (vault_config / "dst.md").read_text()
        assert "New content." in merged
        assert "# Extra" in merged
        assert result["action"] == "content_merged"
        assert result["blocks_added"] == 1

    def test_delete_source_false(self, vault_config):
        """delete_source=False preserves source file."""
        content = "# Same\n"
        (vault_config / "src.md").write_text(content)
        (vault_config / "dst.md").write_text(content)
        result = json.loads(merge_files("src.md", "dst.md", delete_source=False))
        assert result["success"] is True
        assert (vault_config / "src.md").exists()

    def test_source_not_found(self, vault_config):
        (vault_config / "dst.md").write_text("# Dest\n")
        result = json.loads(merge_files("nonexistent.md", "dst.md"))
        assert result["success"] is False
        assert "not found" in result["error"].lower()

    def test_destination_not_found(self, vault_config):
        (vault_config / "src.md").write_text("# Source\n")
        result = json.loads(merge_files("src.md", "nonexistent.md"))
        assert result["success"] is False
        assert "not found" in result["error"].lower()

    def test_concat_strategy(self, vault_config):
        """Concat strategy concatenates without dedup."""
        (vault_config / "src.md").write_text("# Source\n\nSource body.\n")
        (vault_config / "dst.md").write_text("# Dest\n\nDest body.\n")
        result = json.loads(merge_files("src.md", "dst.md", strategy="concat"))
        assert result["success"] is True
        merged = (vault_config / "dst.md").read_text()
        assert "Dest body." in merged
        assert "Source body." in merged
        # Source kept by default for concat
        assert (vault_config / "src.md").exists()

    def test_frontmatter_list_merge(self, vault_config):
        """List fields in frontmatter are unioned."""
        src = "---\ntags:\n  - a\n  - b\n---\n\n# Body\n"
        dst = "---\ntags:\n  - b\n  - c\n---\n\n# Body\n"
        (vault_config / "src.md").write_text(src)
        (vault_config / "dst.md").write_text(dst)
        merge_files("src.md", "dst.md")
        merged = (vault_config / "dst.md").read_text()
        assert "- b" in merged
        assert "- c" in merged
        assert "- a" in merged

    def test_merge_with_heading_positioning(self, vault_config):
        """Unique source block placed after matching heading section in dest."""
        src = "# Log\n\nEntry from source.\n\n# Tasks\n\n- Source task\n"
        dst = "# Tasks\n\n- Dest task\n\n# Log\n\nEntry from dest.\n"
        (vault_config / "src.md").write_text(src)
        (vault_config / "dst.md").write_text(dst)
        result = json.loads(merge_files("src.md", "dst.md"))
        merged = (vault_config / "dst.md").read_text()
        log_pos = merged.index("Entry from dest.")
        source_log_pos = merged.index("Entry from source.")
        tasks_pos = merged.index("- Dest task")
        source_tasks_pos = merged.index("- Source task")
        assert source_log_pos > log_pos
        assert source_tasks_pos > tasks_pos

    def test_invalid_strategy(self, vault_config):
        (vault_config / "src.md").write_text("# A\n")
        (vault_config / "dst.md").write_text("# B\n")
        result = json.loads(merge_files("src.md", "dst.md", strategy="invalid"))
        assert result["success"] is False
        assert "strategy" in result["error"].lower()

    def test_self_merge_rejected(self, vault_config):
        """Merging a file into itself should return error, not delete it."""
        (vault_config / "same.md").write_text("# Content\n")
        result = json.loads(merge_files("same.md", "same.md"))
        assert result["success"] is False
        assert "same file" in result["error"].lower()
        assert (vault_config / "same.md").exists()


class TestBatchMergeFiles:
    """Tests for batch_merge_files tool."""

    def _setup_folders(self, vault_config, pairs):
        """Create source/target folders with file pairs.

        pairs: list of (filename, source_content, target_content)
        """
        src_dir = vault_config / "import"
        dst_dir = vault_config / "Daily Notes"
        src_dir.mkdir(exist_ok=True)
        dst_dir.mkdir(exist_ok=True)
        for name, src_content, dst_content in pairs:
            (src_dir / name).write_text(src_content)
            (dst_dir / name).write_text(dst_content)
        return "import", "Daily Notes"

    def test_batch_merge_identical_files(self, vault_config):
        """All identical pairs: sources deleted, dests unchanged."""
        clear_pending_previews()
        src_folder, dst_folder = self._setup_folders(vault_config, [
            ("2022-01-01.md", "# Jan 1\n\nContent.\n", "# Jan 1\n\nContent.\n"),
            ("2022-01-02.md", "# Jan 2\n\nContent.\n", "# Jan 2\n\nContent.\n"),
        ])
        result = json.loads(batch_merge_files(src_folder, dst_folder))
        assert result["success"] is True
        assert result["merged"] == 2
        assert not (vault_config / "import" / "2022-01-01.md").exists()
        assert not (vault_config / "import" / "2022-01-02.md").exists()

    def test_batch_merge_with_diffs(self, vault_config):
        """Mixed: one identical, one with unique content."""
        clear_pending_previews()
        src_folder, dst_folder = self._setup_folders(vault_config, [
            ("same.md", "# Same\n", "# Same\n"),
            ("diff.md", "# Diff\n\nExtra.\n", "# Diff\n"),
        ])
        result = json.loads(batch_merge_files(src_folder, dst_folder))
        assert result["success"] is True
        assert result["merged"] == 2
        merged = (vault_config / "Daily Notes" / "diff.md").read_text()
        assert "Extra." in merged

    def test_batch_confirmation_gate(self, vault_config):
        """Should require confirmation when >5 pairs."""
        clear_pending_previews()
        src_dir = vault_config / "bulk_src"
        dst_dir = vault_config / "bulk_dst"
        src_dir.mkdir()
        dst_dir.mkdir()
        for i in range(8):
            (src_dir / f"note{i}.md").write_text(f"# Note {i}\n")
            (dst_dir / f"note{i}.md").write_text(f"# Note {i}\n")

        result = json.loads(batch_merge_files("bulk_src", "bulk_dst"))
        assert result["success"] is True
        assert result["confirmation_required"] is True
        assert "8" in result["preview_message"]
        # No files should be merged yet
        for i in range(8):
            assert (src_dir / f"note{i}.md").exists()

    def test_batch_confirm_executes(self, vault_config):
        """Preview then confirm should execute the batch."""
        clear_pending_previews()
        src_dir = vault_config / "conf_src"
        dst_dir = vault_config / "conf_dst"
        src_dir.mkdir()
        dst_dir.mkdir()
        for i in range(8):
            (src_dir / f"n{i}.md").write_text(f"# N {i}\n")
            (dst_dir / f"n{i}.md").write_text(f"# N {i}\n")

        # Preview
        batch_merge_files("conf_src", "conf_dst")
        # Confirm
        result = json.loads(batch_merge_files("conf_src", "conf_dst", confirm=True))
        assert result["success"] is True
        assert result["merged"] == 8

    def test_batch_confirm_rejected_with_different_options(self, vault_config):
        """Changing strategy between preview and confirm requires new preview."""
        clear_pending_previews()
        src_dir = vault_config / "opt_src"
        dst_dir = vault_config / "opt_dst"
        src_dir.mkdir()
        dst_dir.mkdir()
        for i in range(8):
            (src_dir / f"n{i}.md").write_text(f"# N {i}\n")
            (dst_dir / f"n{i}.md").write_text(f"# N {i}\n")

        # Preview with smart strategy
        result = json.loads(batch_merge_files("opt_src", "opt_dst", strategy="smart"))
        assert result["confirmation_required"] is True
        # Confirm with concat strategy — should get a new preview, not execute
        result = json.loads(batch_merge_files("opt_src", "opt_dst", strategy="concat", confirm=True))
        assert result["confirmation_required"] is True
        # All source files should still exist
        for i in range(8):
            assert (src_dir / f"n{i}.md").exists()

    def test_batch_reports_only_in_source(self, vault_config):
        """Files only in source should be reported but not touched."""
        clear_pending_previews()
        src_dir = vault_config / "report_src"
        dst_dir = vault_config / "report_dst"
        src_dir.mkdir()
        dst_dir.mkdir()
        (src_dir / "shared.md").write_text("# Shared\n")
        (dst_dir / "shared.md").write_text("# Shared\n")
        (src_dir / "orphan.md").write_text("# Orphan\n")

        result = json.loads(batch_merge_files("report_src", "report_dst"))
        assert result["success"] is True
        assert result["merged"] == 1
        assert result["skipped_source_only"] == 1
        assert (src_dir / "orphan.md").exists()  # untouched

    def test_batch_no_overlap(self, vault_config):
        """No overlapping files: nothing to merge."""
        clear_pending_previews()
        src_dir = vault_config / "no_overlap_src"
        dst_dir = vault_config / "no_overlap_dst"
        src_dir.mkdir()
        dst_dir.mkdir()
        (src_dir / "a.md").write_text("# A\n")
        (dst_dir / "b.md").write_text("# B\n")

        result = json.loads(batch_merge_files("no_overlap_src", "no_overlap_dst"))
        assert result["success"] is True
        assert result["merged"] == 0

    def test_batch_recursive(self, vault_config):
        """Recursive mode merges files in subfolders."""
        clear_pending_previews()
        src_dir = vault_config / "rec_src"
        dst_dir = vault_config / "rec_dst"
        src_dir.mkdir()
        dst_dir.mkdir()
        (src_dir / "sub").mkdir()
        (src_dir / "sub" / "deep.md").write_text("# Deep\n")
        (dst_dir / "deep.md").write_text("# Deep\n")

        # Non-recursive: no match
        result = json.loads(batch_merge_files("rec_src", "rec_dst"))
        assert result["merged"] == 0

        # Recursive: matches
        result = json.loads(batch_merge_files("rec_src", "rec_dst", recursive=True))
        assert result["merged"] == 1

    def test_batch_skips_ambiguous_targets(self, vault_config):
        """Stems with multiple targets should be skipped, not silently picked."""
        clear_pending_previews()
        src_dir = vault_config / "amb_src"
        dst_dir = vault_config / "amb_dst"
        src_dir.mkdir()
        dst_dir.mkdir()
        (src_dir / "note.md").write_text("# Source\n")
        # Two targets with same stem in different subfolders
        (dst_dir / "sub1").mkdir()
        (dst_dir / "sub2").mkdir()
        (dst_dir / "sub1" / "note.md").write_text("# Target 1\n")
        (dst_dir / "sub2" / "note.md").write_text("# Target 2\n")

        result = json.loads(batch_merge_files("amb_src", "amb_dst", recursive=True))
        assert result["success"] is True
        assert result["merged"] == 0
        assert "note.md" in result["skipped_ambiguous"]
        # Source untouched
        assert (src_dir / "note.md").exists()


class TestReadFileOffice:
    """Tests for read_file Office document dispatch."""

    # --- Word (.docx) ---

    def test_docx_basic(self, vault_config):
        """Should extract heading and paragraph text from a Word document."""
        doc = DocxDocument()
        doc.add_heading("My Title", level=1)
        doc.add_paragraph("First paragraph.")
        doc.add_paragraph("Second paragraph.")
        doc.save(str(vault_config / "test.docx"))

        result = json.loads(read_file("test.docx"))
        assert result["success"] is True
        assert "# My Title" in result["content"]
        assert "First paragraph." in result["content"]
        assert "Second paragraph." in result["content"]

    def test_docx_empty(self, vault_config):
        """Empty Word document should return ok with empty/minimal content."""
        doc = DocxDocument()
        doc.save(str(vault_config / "empty.docx"))

        result = json.loads(read_file("empty.docx"))
        assert result["success"] is True

    def test_docx_with_table(self, vault_config):
        """Should extract table content as markdown table."""
        doc = DocxDocument()
        doc.add_paragraph("Before table.")
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Header A"
        table.cell(0, 1).text = "Header B"
        table.cell(1, 0).text = "Cell 1"
        table.cell(1, 1).text = "Cell 2"
        doc.save(str(vault_config / "table.docx"))

        result = json.loads(read_file("table.docx"))
        assert result["success"] is True
        content = result["content"]
        assert "Header A" in content
        assert "Header B" in content
        assert "Cell 1" in content
        assert "Cell 2" in content
        assert "|" in content  # markdown table syntax
        assert "---" in content  # header separator

    # --- Excel (.xlsx) ---

    def test_xlsx_basic(self, vault_config):
        """Should extract sheet data as markdown table with sheet heading."""
        wb = Workbook()
        ws = wb.active
        ws.title = "Data"
        ws.append(["Name", "Value"])
        ws.append(["Alice", 42])
        wb.save(str(vault_config / "test.xlsx"))

        result = json.loads(read_file("test.xlsx"))
        assert result["success"] is True
        content = result["content"]
        assert "## Data" in content
        assert "Alice" in content
        assert "42" in content
        assert "|" in content
        assert "---" in content

    def test_xlsx_multiple_sheets(self, vault_config):
        """Each sheet should get its own heading."""
        wb = Workbook()
        ws1 = wb.active
        ws1.title = "Sheet1"
        ws1.append(["A", "B"])

        ws2 = wb.create_sheet("Sheet2")
        ws2.append(["C", "D"])
        wb.save(str(vault_config / "multi.xlsx"))

        result = json.loads(read_file("multi.xlsx"))
        assert result["success"] is True
        content = result["content"]
        assert "## Sheet1" in content
        assert "## Sheet2" in content

    def test_xlsx_empty(self, vault_config):
        """Empty workbook should return ok."""
        wb = Workbook()
        # Remove default sheet data (leave it empty)
        ws = wb.active
        # Don't add any data
        wb.save(str(vault_config / "empty.xlsx"))

        result = json.loads(read_file("empty.xlsx"))
        assert result["success"] is True

    # --- PowerPoint (.pptx) ---

    def test_pptx_basic(self, vault_config):
        """Should extract slide title and content."""
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]  # title + content layout
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = "Slide Title"
        slide.placeholders[1].text = "Slide body text."
        prs.save(str(vault_config / "test.pptx"))

        result = json.loads(read_file("test.pptx"))
        assert result["success"] is True
        content = result["content"]
        assert "## Slide Title" in content
        assert "Slide body text." in content

    def test_pptx_multiple_slides(self, vault_config):
        """Multiple slides each get their own heading."""
        prs = Presentation()
        for title_text in ["First", "Second"]:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title_text
            slide.placeholders[1].text = f"Content of {title_text}."
        prs.save(str(vault_config / "multi.pptx"))

        result = json.loads(read_file("multi.pptx"))
        assert result["success"] is True
        content = result["content"]
        assert "## First" in content
        assert "## Second" in content
        assert "Content of First." in content
        assert "Content of Second." in content

    def test_pptx_empty(self, vault_config):
        """Empty presentation should return ok."""
        prs = Presentation()
        prs.save(str(vault_config / "empty.pptx"))

        result = json.loads(read_file("empty.pptx"))
        assert result["success"] is True


class TestExtractBlock:
    """Tests for _extract_block helper."""

    def test_simple_block_id(self):
        """Finds a line with ^blockid and returns it (suffix stripped)."""
        lines = ["# Heading", "- Item one ^abc123", "- Item two"]
        result = _extract_block(lines, "abc123")
        assert result == "- Item one"

    def test_block_with_indented_children(self):
        """Returns the anchor line plus all indented children."""
        lines = [
            "- Parent ^myblock",
            "  - Child 1",
            "  - Child 2",
            "    - Grandchild",
            "- Sibling (not included)",
        ]
        result = _extract_block(lines, "myblock")
        assert result == "- Parent\n  - Child 1\n  - Child 2\n    - Grandchild"

    def test_block_at_end_of_file(self):
        """Block at end of file with children up to EOF."""
        lines = [
            "Some intro",
            "- Last item ^endblock",
            "  - Sub-item",
        ]
        result = _extract_block(lines, "endblock")
        assert result == "- Last item\n  - Sub-item"

    def test_block_not_found(self):
        """Returns None when block ID doesn't exist."""
        lines = ["# Heading", "No blocks here"]
        result = _extract_block(lines, "nonexistent")
        assert result is None

    def test_block_id_mid_line(self):
        """Block ID must be at end of line (after space)."""
        lines = ["Text ^abc123 more text"]
        result = _extract_block(lines, "abc123")
        # Not at end of line, should not match
        assert result is None

    def test_block_no_children(self):
        """Block with no indented children returns just the anchor."""
        lines = [
            "- Item A ^solo",
            "- Item B",
        ]
        result = _extract_block(lines, "solo")
        assert result == "- Item A"


class TestExpandEmbeds:
    """Tests for _expand_embeds — inline embed expansion."""

    def test_no_embeds_unchanged(self, vault_config):
        """Content without embeds is returned unchanged."""
        content = "# Hello\n\nNo embeds here."
        source = vault_config / "source.md"
        result = _expand_embeds(content, source)
        assert result == content

    def test_markdown_full_note_embed(self, vault_config):
        """![[note3]] expands to full note body (no frontmatter)."""
        content = "# Parent\n\n![[note3]]\n\nAfter."
        source = vault_config / "parent.md"
        result = _expand_embeds(content, source)
        assert "> [Embedded: note3]" in result
        assert "> # Note 3" in result
        assert "![[note3]]" not in result
        assert "After." in result

    def test_markdown_embed_strips_frontmatter(self, vault_config):
        """Embedded markdown notes have frontmatter stripped."""
        content = "Before\n\n![[note1]]\n\nAfter"
        source = vault_config / "parent.md"
        result = _expand_embeds(content, source)
        assert "> [Embedded: note1]" in result
        # Frontmatter delimiters should not appear in the embedded section
        embedded_section = result.split("> [Embedded: note1]")[1].split("After")[0]
        assert "---" not in embedded_section
        assert "> # Note 1" in result

    def test_heading_embed(self, vault_config):
        """![[note2#Section A]] expands only that section."""
        content = "See: ![[note2#Section A]]"
        source = vault_config / "parent.md"
        result = _expand_embeds(content, source)
        assert "> [Embedded: note2#Section A]" in result
        assert "Content in section A" in result
        assert "Content in section B" not in result

    def test_block_id_embed(self, vault_config):
        """![[note#^blockid]] expands the block and its children."""
        (vault_config / "blocks.md").write_text(
            "# Blocks\n\n- Item one ^myid\n  - Child\n- Other\n"
        )
        content = "Reference: ![[blocks#^myid]]"
        source = vault_config / "parent.md"
        result = _expand_embeds(content, source)
        assert "> [Embedded: blocks#^myid]" in result
        assert "Item one" in result
        assert "Child" in result
        assert "Other" not in result

    def test_unresolved_embed_error_marker(self, vault_config):
        """Unresolvable embeds produce an error marker."""
        content = "![[nonexistent_file]]"
        source = vault_config / "parent.md"
        result = _expand_embeds(content, source)
        assert "> [Embed error: nonexistent_file" in result

    def test_self_embed_skipped(self, vault_config):
        """Self-referencing embeds produce an error marker."""
        (vault_config / "self.md").write_text("# Self\n\n![[self]]\n")
        result = _expand_embeds("# Self\n\n![[self]]\n", vault_config / "self.md")
        assert "> [Embed error: self" in result
        assert "self-reference" in result.lower()

    def test_embed_in_code_block_not_expanded(self, vault_config):
        """Embeds inside fenced code blocks are left as-is."""
        content = "# Doc\n\n```\n![[note3]]\n```\n\n![[note3]]\n"
        source = vault_config / "parent.md"
        result = _expand_embeds(content, source)
        # The one inside the code block should be literal
        assert "```\n![[note3]]\n```" in result
        # The one outside should be expanded
        assert "> [Embedded: note3]" in result

    def test_embed_in_inline_code_not_expanded(self, vault_config):
        """Embeds inside inline code backticks are left as-is."""
        content = "Use `![[note3]]` to embed files."
        source = vault_config / "parent.md"
        result = _expand_embeds(content, source)
        assert "![[note3]]" in result
        assert "> [Embedded:" not in result

    def test_embed_with_dot_in_folder_name(self, vault_config):
        """![[2026.02/daily]] resolves as markdown despite dot in folder name."""
        folder = vault_config / "2026.02"
        folder.mkdir()
        (folder / "daily.md").write_text("# Daily\n\nToday's notes.\n")
        content = "![[2026.02/daily]]"
        source = vault_config / "parent.md"
        result = _expand_embeds(content, source)
        assert "> [Embedded: 2026.02/daily]" in result
        assert "Today's notes" in result

    def test_aliased_embed(self, vault_config):
        """![[note3|Summary]] strips alias and expands the note."""
        content = "![[note3|Summary]]"
        source = vault_config / "parent.md"
        result = _expand_embeds(content, source)
        assert "> [Embedded: note3|Summary]" in result
        assert "> # Note 3" in result

    def test_aliased_heading_embed(self, vault_config):
        """![[note2#Section A|see here]] strips alias and expands section."""
        content = "![[note2#Section A|see here]]"
        source = vault_config / "parent.md"
        result = _expand_embeds(content, source)
        assert "Content in section A" in result
        assert "Content in section B" not in result

    def test_multiple_embeds(self, vault_config):
        """Multiple embeds in one file are all expanded."""
        content = "![[note1]]\n\n![[note3]]"
        source = vault_config / "parent.md"
        result = _expand_embeds(content, source)
        assert "> [Embedded: note1]" in result
        assert "> [Embedded: note3]" in result

    def test_binary_embed_audio(self, vault_config, monkeypatch):
        """Audio embeds call handle_audio and format the result."""
        monkeypatch.setenv("FIREWORKS_API_KEY", "test-key")
        audio = vault_config / "Attachments" / "rec.m4a"
        audio.write_bytes(b"fake audio")

        from unittest.mock import patch as _patch
        with _patch("tools.files.handle_audio") as mock_audio:
            mock_audio.return_value = '{"success": true, "transcript": "Hello world"}'
            content = "![[rec.m4a]]"
            source = vault_config / "parent.md"
            _embed_cache.clear()
            result = _expand_embeds(content, source)
            assert "> [Embedded: rec.m4a]" in result
            assert "> Hello world" in result

    def test_binary_embed_cache_hit(self, vault_config, monkeypatch):
        """Second expansion of same binary embed uses cache."""
        monkeypatch.setenv("FIREWORKS_API_KEY", "test-key")
        audio = vault_config / "Attachments" / "rec.m4a"
        audio.write_bytes(b"fake audio")

        from unittest.mock import patch as _patch
        with _patch("tools.files.handle_audio") as mock_audio:
            mock_audio.return_value = '{"success": true, "transcript": "Cached"}'
            content = "![[rec.m4a]]"
            source = vault_config / "parent.md"
            _embed_cache.clear()
            _expand_embeds(content, source)
            _expand_embeds(content, source)
            assert mock_audio.call_count == 1

    def test_binary_embed_cache_miss_logged(self, vault_config, monkeypatch, caplog):
        """Cache miss is logged at DEBUG level with filename and handler type."""
        monkeypatch.setenv("FIREWORKS_API_KEY", "test-key")
        audio = vault_config / "Attachments" / "rec.m4a"
        audio.write_bytes(b"fake audio")

        with patch("tools.files.handle_audio") as mock_audio:
            mock_audio.return_value = '{"success": true, "transcript": "Hello"}'
            _embed_cache.clear()
            with caplog.at_level(logging.DEBUG, logger="tools.files"):
                _expand_embeds("![[rec.m4a]]", vault_config / "parent.md")
            assert any("Cache miss" in r.message and "rec.m4a" in r.message
                       for r in caplog.records)

    def test_binary_embed_cache_hit_logged(self, vault_config, monkeypatch, caplog):
        """Cache hit is logged at DEBUG level with filename."""
        monkeypatch.setenv("FIREWORKS_API_KEY", "test-key")
        audio = vault_config / "Attachments" / "rec.m4a"
        audio.write_bytes(b"fake audio")

        with patch("tools.files.handle_audio") as mock_audio:
            mock_audio.return_value = '{"success": true, "transcript": "Hello"}'
            _embed_cache.clear()
            _expand_embeds("![[rec.m4a]]", vault_config / "parent.md")
            with caplog.at_level(logging.DEBUG, logger="tools.files"):
                _expand_embeds("![[rec.m4a]]", vault_config / "parent.md")
            assert any("Cache hit" in r.message and "rec.m4a" in r.message
                       for r in caplog.records)


class TestBinaryHandlerLogging:
    """Tests for logging in binary embed handlers."""

    def test_handle_audio_logs_entry_and_success(self, tmp_path, caplog):
        """handle_audio logs file name, size, and duration on success."""
        audio = tmp_path / "rec.m4a"
        audio.write_bytes(b"x" * 1024)

        mock_response = MagicMock()
        mock_response.text = "Hello world"

        with patch("tools.readers.OpenAI") as mock_openai_cls:
            mock_client = MagicMock()
            mock_openai_cls.return_value = mock_client
            mock_client.audio.transcriptions.create.return_value = mock_response
            with patch.dict("os.environ", {"FIREWORKS_API_KEY": "test-key"}):
                with caplog.at_level(logging.INFO, logger="tools.readers"):
                    result = handle_audio(audio)

        assert json.loads(result)["success"] is True
        messages = [r.message for r in caplog.records]
        assert any("rec.m4a" in m and "1024" in m for m in messages), \
            f"Expected entry log with filename and size, got: {messages}"
        assert any("Transcribed" in m and "rec.m4a" in m for m in messages), \
            f"Expected success log with 'Transcribed' and filename, got: {messages}"

    def test_handle_audio_logs_warning_on_failure(self, tmp_path, caplog):
        """handle_audio logs a WARNING when the API call raises."""
        audio = tmp_path / "bad.m4a"
        audio.write_bytes(b"data")

        with patch("tools.readers.OpenAI") as mock_openai_cls:
            mock_client = MagicMock()
            mock_openai_cls.return_value = mock_client
            mock_client.audio.transcriptions.create.side_effect = RuntimeError("API down")
            with patch.dict("os.environ", {"FIREWORKS_API_KEY": "test-key"}):
                with caplog.at_level(logging.WARNING, logger="tools.readers"):
                    handle_audio(audio)

        assert any("bad.m4a" in r.message and r.levelname == "WARNING"
                   for r in caplog.records)

    def test_handle_image_logs_entry_and_success(self, tmp_path, caplog):
        """handle_image logs file name, size, and duration on success."""
        img = tmp_path / "diagram.png"
        img.write_bytes(b"x" * 2048)

        mock_response = MagicMock()
        mock_response.choices = [MagicMock()]
        mock_response.choices[0].message.content = "A diagram"

        with patch("tools.readers.OpenAI") as mock_openai_cls:
            mock_client = MagicMock()
            mock_openai_cls.return_value = mock_client
            mock_client.chat.completions.create.return_value = mock_response
            with patch.dict("os.environ", {"FIREWORKS_API_KEY": "test-key"}):
                with caplog.at_level(logging.INFO, logger="tools.readers"):
                    result = handle_image(img)

        assert json.loads(result)["success"] is True
        messages = [r.message for r in caplog.records]
        assert any("diagram.png" in m and "2048" in m for m in messages), \
            f"Expected entry log with filename and size, got: {messages}"
        assert any("Described" in m and "diagram.png" in m for m in messages), \
            f"Expected success log with 'Described' and filename, got: {messages}"

    def test_handle_image_logs_warning_on_failure(self, tmp_path, caplog):
        """handle_image logs a WARNING when the API call raises."""
        img = tmp_path / "broken.png"
        img.write_bytes(b"data")

        with patch("tools.readers.OpenAI") as mock_openai_cls:
            mock_client = MagicMock()
            mock_openai_cls.return_value = mock_client
            mock_client.chat.completions.create.side_effect = RuntimeError("timeout")
            with patch.dict("os.environ", {"FIREWORKS_API_KEY": "test-key"}):
                with caplog.at_level(logging.WARNING, logger="tools.readers"):
                    result = handle_image(img)

        assert json.loads(result)["success"] is False
        assert any("broken.png" in r.message and r.levelname == "WARNING"
                   for r in caplog.records)

    def test_handle_office_logs_entry_and_success(self, tmp_path, caplog):
        """handle_office logs file name, size, and duration on success."""
        from docx import Document as DocxDocument
        docx = tmp_path / "report.docx"
        doc = DocxDocument()
        doc.add_paragraph("Hello")
        doc.save(str(docx))
        size = docx.stat().st_size

        with caplog.at_level(logging.INFO, logger="tools.readers"):
            result = handle_office(docx)

        assert json.loads(result)["success"] is True
        messages = [r.message for r in caplog.records]
        assert any("report.docx" in m and str(size) in m for m in messages), \
            f"Expected entry log with filename and size, got: {messages}"
        assert any("Extracted" in m and "report.docx" in m for m in messages), \
            f"Expected success log with 'Extracted' and filename, got: {messages}"

    def test_handle_office_logs_warning_on_failure(self, tmp_path, caplog):
        """handle_office logs a WARNING when extraction fails."""
        bad = tmp_path / "corrupt.docx"
        bad.write_bytes(b"not a real docx")

        with caplog.at_level(logging.WARNING, logger="tools.readers"):
            result = handle_office(bad)

        assert json.loads(result)["success"] is False
        assert any("corrupt.docx" in r.message and r.levelname == "WARNING"
                   for r in caplog.records)


class TestExtractHeadings:
    """Tests for _extract_headings helper."""

    def test_basic_headings(self):
        """Should extract headings from markdown content."""
        content = "# Title\n\nSome text\n\n## Section 1\n\nMore text\n\n### Subsection\n\n## Section 2\n"
        assert _extract_headings(content) == ["# Title", "## Section 1", "### Subsection", "## Section 2"]

    def test_no_headings(self):
        """Should return empty list for content without headings."""
        assert _extract_headings("Just plain text\nwith lines\n") == []

    def test_headings_inside_code_fence_skipped(self):
        """Should skip headings inside code fences."""
        content = "# Real Heading\n\n```\n# Not a heading\n## Also not\n```\n\n## Real Section\n"
        assert _extract_headings(content) == ["# Real Heading", "## Real Section"]

    def test_tilde_code_fence(self):
        """Should skip headings inside tilde fences."""
        content = "# Title\n\n~~~\n## Fake\n~~~\n\n## Real\n"
        assert _extract_headings(content) == ["# Title", "## Real"]

    def test_empty_content(self):
        """Should return empty list for empty string."""
        assert _extract_headings("") == []

    def test_mismatched_fence_delimiters(self):
        """~~~ inside a ``` block should not toggle fence state."""
        content = "# Before\n\n```\n~~~\n## Fake\n~~~\n```\n\n## After\n"
        assert _extract_headings(content) == ["# Before", "## After"]

    def test_backticks_inside_tilde_fence(self):
        """``` inside a ~~~ block should not toggle fence state."""
        content = "# Before\n\n~~~\n```\n## Fake\n```\n~~~\n\n## After\n"
        assert _extract_headings(content) == ["# Before", "## After"]

    def test_frontmatter_comments_excluded(self):
        """YAML comments (# ...) in frontmatter should not appear as headings."""
        content = "---\n# This is a YAML comment\ntags:\n  - test\n---\n\n## Real Heading\n"
        assert _extract_headings(content) == ["## Real Heading"]

    def test_longer_fence_not_closed_by_shorter(self):
        """A ```` block should not be closed by ```."""
        content = "# Before\n\n````\n```\n## Fake\n```\n````\n\n## After\n"
        assert _extract_headings(content) == ["# Before", "## After"]

    def test_longer_tilde_fence_not_closed_by_shorter(self):
        """A ~~~~ block should not be closed by ~~~."""
        content = "# Before\n\n~~~~\n~~~\n## Fake\n~~~\n~~~~\n\n## After\n"
        assert _extract_headings(content) == ["# Before", "## After"]


class TestGetNoteInfo:
    """Tests for get_note_info tool."""

    def test_basic_metadata(self, vault_config, temp_vault):
        """Should return frontmatter, headings, size, timestamps."""
        result = json.loads(get_note_info("note1.md"))
        assert result["success"] is True
        assert result["path"] == "note1.md"
        assert result["frontmatter"]["tags"] == ["project", "work"]
        assert "# Note 1" in result["headings"]
        assert isinstance(result["size"], int)
        assert result["size"] > 0
        assert "modified" in result
        assert "created" in result

    def test_link_counts(self, vault_config, temp_vault):
        """Should include backlink and outlink counts."""
        result = json.loads(get_note_info("note1.md"))
        assert "backlink_count" in result
        assert "outlink_count" in result
        assert isinstance(result["backlink_count"], int)
        assert isinstance(result["outlink_count"], int)
        # note1 has [[wikilink]] outlink
        assert result["outlink_count"] >= 1

    def test_no_frontmatter(self, vault_config, temp_vault):
        """Should return empty frontmatter for files without it."""
        (temp_vault / "plain.md").write_text("# Just a heading\n\nNo frontmatter here.\n")
        result = json.loads(get_note_info("plain.md"))
        assert result["success"] is True
        assert result["frontmatter"] == {}
        assert "# Just a heading" in result["headings"]

    def test_nonexistent_file(self, vault_config):
        """Should return error for missing file."""
        result = json.loads(get_note_info("nonexistent.md"))
        assert result["success"] is False
        assert "not found" in result["error"].lower()

    def test_non_markdown_file(self, vault_config, temp_vault):
        """Should return basic metadata for non-markdown files."""
        csv_file = temp_vault / "data.csv"
        csv_file.write_text("a,b,c\n1,2,3\n")
        result = json.loads(get_note_info("data.csv"))
        assert result["success"] is True
        assert result["frontmatter"] == {}
        assert result["headings"] == []
        assert result["backlink_count"] == 0
        assert result["outlink_count"] == 0
        # Non-markdown uses stat byte size, not decoded char count
        assert result["size"] == csv_file.stat().st_size

    def test_headings_respect_code_fences(self, vault_config, temp_vault):
        """Should skip headings inside code fences."""
        (temp_vault / "fenced.md").write_text(
            "# Real\n\n```\n## Fake\n```\n\n## Also Real\n"
        )
        result = json.loads(get_note_info("fenced.md"))
        assert result["headings"] == ["# Real", "## Also Real"]

    def test_created_from_frontmatter_date(self, vault_config, temp_vault):
        """Should use frontmatter Date field for created timestamp."""
        result = json.loads(get_note_info("note1.md"))
        # note1.md has Date: 2024-01-15
        assert result["created"].startswith("2024-01-15")

    def test_nbs_in_path(self, vault_config, temp_vault):
        """Should handle non-breaking spaces in path."""
        (temp_vault / "test nbs.md").write_text("# Test\n")
        result = json.loads(get_note_info("test\xa0nbs.md"))
        assert result["success"] is True

    def test_attachment_fallback(self, vault_config, temp_vault):
        """Should resolve bare binary filenames via Attachments directory."""
        att_dir = temp_vault / "Attachments"
        att_dir.mkdir(exist_ok=True)
        img = att_dir / "diagram.png"
        img.write_bytes(b"\x89PNG\r\n\x1a\n" + b"\x00" * 100)
        result = json.loads(get_note_info("diagram.png"))
        assert result["success"] is True
        assert result["size"] == img.stat().st_size

    def test_non_dict_frontmatter(self, vault_config, temp_vault):
        """Should return empty frontmatter when YAML parses to non-dict."""
        (temp_vault / "bad_fm.md").write_text("---\n- item1\n- item2\n---\n\n# Heading\n")
        result = json.loads(get_note_info("bad_fm.md"))
        assert result["success"] is True
        assert result["frontmatter"] == {}
        assert result["headings"] == ["# Heading"]

    def test_frontmatter_no_trailing_newline(self, vault_config, temp_vault):
        """Should parse frontmatter when file ends without trailing newline after ---."""
        (temp_vault / "no_nl.md").write_text("---\nDate: 2024-06-15\ntags:\n  - test\n---")
        result = json.loads(get_note_info("no_nl.md"))
        assert result["success"] is True
        assert result["frontmatter"]["tags"] == ["test"]
        assert result["created"].startswith("2024-06-15")
