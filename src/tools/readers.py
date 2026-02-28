"""File type handlers for read_file dispatch.

Each handler takes a resolved Path and returns ok()/err() JSON.
"""

import base64
import logging
import os
import time
from pathlib import Path

from openai import OpenAI

from config import FIREWORKS_BASE_URL, VISION_MODEL, WHISPER_MODEL
from services.vault import err, ok

logger = logging.getLogger(__name__)

# Extension sets for dispatch
AUDIO_EXTENSIONS = {".m4a", ".mp3", ".wav", ".ogg", ".webm"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".svg"}
OFFICE_EXTENSIONS = {".docx", ".xlsx", ".pptx"}
PDF_EXTENSIONS = {".pdf"}


def handle_audio(file_path: Path) -> str:
    """Transcribe an audio file using Fireworks Whisper API."""
    api_key = os.getenv("FIREWORKS_API_KEY")
    if not api_key:
        return err("FIREWORKS_API_KEY not set")

    try:
        size = file_path.stat().st_size
    except OSError:
        size = 0

    logger.info("Transcribing audio: %s (%d bytes)", file_path.name, size)
    client = OpenAI(api_key=api_key, base_url=FIREWORKS_BASE_URL)
    start = time.perf_counter()
    try:
        with open(file_path, "rb") as f:
            response = client.audio.transcriptions.create(
                model=WHISPER_MODEL,
                file=f,
                response_format="verbose_json",
                timestamp_granularities=["word"],
                extra_body={"diarize": True},
            )
        elapsed = time.perf_counter() - start
        logger.info("Transcribed %s in %.2fs", file_path.name, elapsed)
        return ok(transcript=response.text)
    except Exception as e:
        logger.warning("Transcription failed for %s: %s", file_path.name, e)
        return err(f"Transcription failed: {e}")


def handle_image(file_path: Path) -> str:
    """Describe an image using Fireworks vision model."""
    api_key = os.getenv("FIREWORKS_API_KEY")
    if not api_key:
        return err("FIREWORKS_API_KEY not set")

    try:
        size = file_path.stat().st_size
    except OSError:
        size = 0

    logger.info("Describing image: %s (%d bytes)", file_path.name, size)
    client = OpenAI(api_key=api_key, base_url=FIREWORKS_BASE_URL)
    start = time.perf_counter()
    try:
        image_data = file_path.read_bytes()
        b64 = base64.b64encode(image_data).decode("utf-8")

        mime_map = {
            ".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
            ".gif": "image/gif", ".webp": "image/webp", ".svg": "image/svg+xml",
        }
        mime = mime_map.get(file_path.suffix.lower(), "image/png")

        response = client.chat.completions.create(
            model=VISION_MODEL,
            messages=[{
                "role": "user",
                "content": [
                    {"type": "text", "text": "Describe this image in detail, including any visible text."},
                    {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64}"}},
                ],
            }],
        )
        elapsed = time.perf_counter() - start
        description = response.choices[0].message.content
        logger.info("Described %s in %.2fs", file_path.name, elapsed)
        return ok(description=description)
    except Exception as e:
        logger.warning("Image description failed for %s: %s", file_path.name, e)
        return err(f"Image description failed: {e}")


def handle_office(file_path: Path) -> str:
    """Extract text content from Office documents (.docx, .xlsx, .pptx)."""
    try:
        size = file_path.stat().st_size
    except OSError:
        size = 0

    logger.info("Extracting office document: %s (%d bytes)", file_path.name, size)
    ext = file_path.suffix.lower()
    start = time.perf_counter()
    try:
        if ext == ".docx":
            result = _read_docx(file_path)
        elif ext == ".xlsx":
            result = _read_xlsx(file_path)
        elif ext == ".pptx":
            result = _read_pptx(file_path)
        else:
            return err(f"Unsupported office format: {ext}")
        elapsed = time.perf_counter() - start
        logger.info("Extracted %s in %.2fs", file_path.name, elapsed)
        return result
    except Exception as e:
        logger.warning("Office extraction failed for %s: %s", file_path.name, e)
        return err(f"Failed to read {file_path.name}: {e}")


def handle_pdf(file_path: Path) -> str:
    """Extract text content from a PDF file, page by page."""
    import pymupdf

    try:
        size = file_path.stat().st_size
    except OSError:
        size = 0

    logger.info("Extracting PDF: %s (%d bytes)", file_path.name, size)
    start = time.perf_counter()
    try:
        with pymupdf.Document(str(file_path)) as doc:
            parts = []
            for i, page in enumerate(doc, 1):
                text = page.get_text().strip()
                if text:
                    parts.append(f"## Page {i}\n\n{text}")
        elapsed = time.perf_counter() - start
        logger.info("Extracted %s in %.2fs", file_path.name, elapsed)
        content = "\n\n".join(parts)
        return ok(content=content)
    except Exception as e:
        logger.warning("PDF extraction failed for %s: %s", file_path.name, e)
        return err(f"Failed to read {file_path.name}: {e}")


def _read_docx(file_path: Path) -> str:
    """Extract text from a Word document, preserving headings and tables."""
    from docx import Document as DocxDocument
    from docx.oxml.ns import qn
    from docx.table import Table

    doc = DocxDocument(str(file_path))
    parts = []

    for element in doc.element.body:
        tag = element.tag.split("}")[-1]  # strip namespace
        if tag == "p":
            # Check if it's a heading
            style = element.find(qn("w:pPr"))
            if style is not None:
                pstyle = style.find(qn("w:pStyle"))
                if pstyle is not None:
                    val = pstyle.get(qn("w:val"), "")
                    if val.startswith("Heading"):
                        level = val.replace("Heading", "").strip()
                        try:
                            level = int(level)
                        except ValueError:
                            level = 1
                        text = element.text or ""
                        if text.strip():
                            parts.append(f"{'#' * level} {text.strip()}")
                            continue
            text = element.text or ""
            if text.strip():
                parts.append(text.strip())
        elif tag == "tbl":
            # Extract table as markdown
            table = Table(element, doc)
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append("| " + " | ".join(cells) + " |")
            if rows:
                header_sep = "| " + " | ".join("---" for _ in table.rows[0].cells) + " |"
                rows.insert(1, header_sep)
                parts.append("\n".join(rows))

    content = "\n\n".join(parts)
    return ok(content=content)


def _read_xlsx(file_path: Path) -> str:
    """Extract Excel data as markdown tables, one per sheet."""
    from openpyxl import load_workbook

    wb = load_workbook(str(file_path), read_only=True, data_only=True)
    parts = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(cell) if cell is not None else "" for cell in row]
            if any(c for c in cells):  # skip fully empty rows
                rows.append(cells)

        if not rows:
            continue

        parts.append(f"## {sheet_name}")
        table_rows = ["| " + " | ".join(cells) + " |" for cells in rows]
        header_sep = "| " + " | ".join("---" for _ in rows[0]) + " |"
        table_rows.insert(1, header_sep)
        parts.append("\n".join(table_rows))

    wb.close()
    content = "\n\n".join(parts)
    return ok(content=content)


def _read_pptx(file_path: Path) -> str:
    """Extract text from PowerPoint slides."""
    from pptx import Presentation

    prs = Presentation(str(file_path))
    parts = []

    for i, slide in enumerate(prs.slides, 1):
        slide_parts = []
        title = None

        # Extract title if present
        if slide.shapes.title and slide.shapes.title.text.strip():
            title = slide.shapes.title.text.strip()

        heading = f"## {title}" if title else f"## Slide {i}"
        slide_parts.append(heading)

        # Extract text from all shapes (excluding title to avoid duplication)
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    slide_parts.append(text)

        if len(slide_parts) > 1:  # has content beyond heading
            parts.append("\n\n".join(slide_parts))
        elif title:  # title-only slide still worth including
            parts.append(heading)

    content = "\n\n".join(parts)
    return ok(content=content)
